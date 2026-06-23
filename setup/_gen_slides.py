
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG      = RGBColor(0x0d, 0x11, 0x17)
GREEN   = RGBColor(0x10, 0xb9, 0x81)
CARD_BG = RGBColor(0x1e, 0x25, 0x35)
WHITE   = RGBColor(0xff, 0xff, 0xff)
GRAY    = RGBColor(0x8b, 0x95, 0xa7)
AMBER   = RGBColor(0xf5, 0x9e, 0x0b)
DARK    = RGBColor(0x0a, 0x0e, 0x14)
BLACK   = RGBColor(0x00, 0x00, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def txb(s, text, x, y, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color

def box(s, x, y, w, h, color=CARD_BG):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.color.rgb = color
    return sh

def divider(s, y):
    box(s, 0.5, y, 12.33, 0.03, GREEN)

def pill(s, text, x, y, bg=GREEN, fg=BLACK):
    w = len(text) * 0.115 + 0.2
    box(s, x, y, w, 0.38, bg)
    txb(s, text, x+0.08, y+0.05, w-0.1, 0.28, size=12, bold=True, color=fg, align=PP_ALIGN.CENTER)

def pilar_header(s, num, emoji, nome, atitulo, atexto):
    pill(s, "Pilar " + str(num), 0.5, 0.28)
    txb(s, emoji + "  " + nome, 1.9, 0.25, 10.5, 0.65, size=28, bold=True)
    divider(s, 1.05)
    box(s, 0.5, 1.2, 12.3, 0.3, GREEN)
    txb(s, "Analogia: " + atitulo, 0.65, 1.22, 11.8, 0.26, size=12, bold=True, color=BLACK)
    box(s, 0.5, 1.5, 12.3, 1.0, CARD_BG)
    txb(s, atexto, 0.65, 1.55, 12.0, 0.9, size=13, italic=True)

# ── SLIDE 1: TITULO ──────────────────────────────────────────
s = slide()
r = s.shapes.add_shape(1, Inches(0), Inches(3.0), Inches(0.12), Inches(1.5))
r.fill.solid(); r.fill.fore_color.rgb = GREEN; r.line.color.rgb = GREEN
txb(s, "GitHub Copilot", 0.4, 1.5, 12, 1.2, size=52, bold=True)
txb(s, "Platform Extensibility", 0.4, 2.7, 12, 1.0, size=36, color=GREEN)
divider(s, 4.0)
txb(s, "Academia FitCode  .  Workshop Fase 2", 0.4, 4.15, 12, 0.5, size=18, color=GRAY)
txb(s, "8 Pilares de Governanca do Copilot", 0.4, 4.75, 12, 0.5, size=16, color=GRAY, italic=True)

# ── SLIDE 2: OBJETIVO ────────────────────────────────────────
s = slide()
txb(s, "Objetivo do Workshop", 0.5, 0.4, 12, 0.8, size=36, bold=True, color=GREEN)
divider(s, 1.35)
txb(s, "Criar um site funcional e configurar os 8 pilares de governanca do GitHub Copilot -- um a um.", 0.5, 1.6, 12, 0.9, size=22)
box(s, 0.5, 2.75, 12.3, 1.3, CARD_BG)
txb(s, "O challenge termina quando voce rodar gerar_evidencia.py\ne ver 8/8 pilares configurados no HTML de evidencia.", 0.7, 2.95, 11.8, 1.0, size=18, color=GREEN)
txb(s, "Antes de comecar:", 0.5, 4.3, 5, 0.4, size=18, bold=True)
txb(s, "Siga o guia PASSOS-INICIAIS.md para configurar seu ambiente.", 0.5, 4.8, 12, 0.4, size=16, color=GRAY, italic=True)

# ── SLIDE 3: 8 PILARES VISAO GERAL ───────────────────────────
s = slide()
txb(s, "Os 8 Pilares", 0.5, 0.3, 10, 0.7, size=36, bold=True, color=GREEN)
divider(s, 1.1)
pilares = [
    ("1","Instructions",  "Regras permanentes que o Copilot\nsegue em todo contexto"),
    ("2","Skills",         "Conhecimento especializado\ninjetado sob demanda"),
    ("3","Agents",         "Especialista com identidade,\nescopo e ferramentas proprias"),
    ("4","MCP",            "Conexao a ferramentas externas\nem tempo real"),
    ("5","Hooks",          "Intercepta e valida acoes\ndo agente automaticamente"),
    ("6","Knowledge Base", "Documentacao versionada\ncomo fonte de verdade"),
    ("7","Coding Agent",   "Copilot autonomo no GitHub\nIssue -> PR automatico"),
    ("8","CLI",            "Copilot diretamente\nno terminal"),
]
cols = [0.3, 3.55, 6.8, 10.05]
rows = [1.3, 3.8]
for i,(num,nome,desc) in enumerate(pilares):
    cx = cols[i%4]; cy = rows[i//4]
    box(s, cx, cy, 3.0, 2.2, CARD_BG)
    r2 = s.shapes.add_shape(1,Inches(cx+0.1),Inches(cy+0.1),Inches(0.42),Inches(0.42))
    r2.fill.solid(); r2.fill.fore_color.rgb = GREEN; r2.line.color.rgb = GREEN
    txb(s, num, cx+0.12, cy+0.1, 0.38, 0.38, size=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    txb(s, nome, cx+0.65, cy+0.12, 2.2, 0.4, size=14, bold=True)
    txb(s, desc, cx+0.12, cy+0.65, 2.76, 0.9, size=11, color=GRAY)

# ── SLIDE 4: PASSO 0 ─────────────────────────────────────────
s = slide()
pill(s, "Passo 0", 0.5, 0.3)
txb(s, "Crie o Site com o Copilot", 1.7, 0.27, 10.5, 0.6, size=28, bold=True)
divider(s, 1.05)
txb(s, "Antes de aplicar governanca, use o Copilot para criar as secoes do site.", 0.5, 1.2, 12.3, 0.6, size=16, color=GRAY)
box(s, 0.5, 2.0, 12.3, 2.5, DARK)
txb(s, 'Crie a secao "Sobre a Academia" do site em HTML semantico.\nInclua: titulo h2, paragrafo de apresentacao, e tres cards\ncom diferenciais (Profissionais certificados, Equipamentos modernos,\nLocalizacao central). Use classes CSS com prefixo "fc-".',
    0.7, 2.15, 11.8, 2.2, size=14, color=GREEN)
txb(s, "ATENCAO: Anote o resultado -- ele provavelmente nao vai seguir padroes de acessibilidade\n(aria-label, roles, contraste) nem SEO. Voce vai ver a diferenca apos configurar as instructions.",
    0.5, 4.7, 12.3, 0.9, size=14, color=AMBER)
txb(s, "Criterio: a secao existe e renderiza no navegador.", 0.5, 5.8, 12.3, 0.4, size=14, color=GREEN)

# ── SLIDE 5: PILAR 1 INSTRUCTIONS ────────────────────────────
s = slide()
pilar_header(s, "1", "INSTRUCOES", "Instructions",
    "Regras do condominio",
    "Voce recebe um regulamento: pode fazer barulho ate 22h, lixo vai na lixeira separada.\nVoce nao precisa perguntar essas regras toda vez -- elas ja foram internalizadas.")
txb(s, "O que e:", 0.5, 2.65, 3, 0.35, size=14, bold=True, color=GREEN)
txb(s, "Regras permanentes que o Copilot le em todo contexto, automaticamente.", 0.5, 3.05, 12.3, 0.4, size=15)
items = [
    ("copilot-instructions.md", "padrao HTML semantico, nomenclatura CSS (fc-), regras de imagem"),
    ("acessibilidade.instructions.md", "aria-label, roles, contraste WCAG AA  --  applyTo: **/*.html"),
    ("seo.instructions.md", "meta tags obrigatorias, Open Graph, alt  --  applyTo: **/*.html"),
]
for i,(f,d) in enumerate(items):
    box(s, 0.5, 3.55+i*0.62, 12.3, 0.52, CARD_BG)
    txb(s, f, 0.7, 3.58+i*0.62, 4.5, 0.28, size=12, bold=True, color=GREEN)
    txb(s, d, 5.4, 3.58+i*0.62, 7.2, 0.28, size=12, color=GRAY)
txb(s, "Valide: peca ao Copilot para criar um card -- deve gerar com aria-label, fc- e alt.", 0.5, 6.7, 12.3, 0.4, size=13, color=GREEN)

# ── SLIDE 6: PILAR 2 SKILLS ──────────────────────────────────
s = slide()
pilar_header(s, "2", "SKILLS", "Skills",
    "Cardapio de especialidades do chef",
    "O chef tem fichas tecnicas especificas: temperatura exata, tempo de cozimento.\nQuando um pedido especial entra, ele puxa a ficha tecnica. As Skills sao essas fichas.")
txb(s, "O que e:", 0.5, 2.65, 3, 0.35, size=14, bold=True, color=GREEN)
txb(s, "Conhecimento especializado injetado sob demanda quando o contexto e relevante.", 0.5, 3.05, 12.3, 0.4, size=15)
headers = ["", "Instructions", "Skills"]
rows2 = [
    ["Quando ativa", "Sempre -- em todo chat do repo", "Sob demanda -- contexto relevante"],
    ["O que define", "Regras gerais (semantica, a11y, SEO)", "Procedimentos especificos (templates)"],
    ["Exemplo", '"Toda imagem deve ter alt"', '"Ao gerar secao hero, use este template"'],
]
cx2 = [0.5, 3.1, 7.7]; cw2 = [2.5, 4.5, 4.5]
box(s, 0.5, 3.55, 11.7, 0.38, GREEN)
for ci,h in enumerate(headers):
    txb(s, h, cx2[ci]+0.1, 3.57, cw2[ci], 0.28, size=12, bold=True, color=BLACK)
for ri,row in enumerate(rows2):
    bg2 = CARD_BG if ri%2==0 else RGBColor(0x16,0x1c,0x2a)
    box(s, 0.5, 3.93+ri*0.65, 11.7, 0.6, bg2)
    for ci,cell in enumerate(row):
        col2 = GREEN if ci==0 else (WHITE if ci==1 else GRAY)
        txb(s, cell, cx2[ci]+0.1, 3.96+ri*0.65, cw2[ci]-0.1, 0.52, size=11, color=col2)
txb(s, "Valide: peca 'Gere a secao de depoimentos' -- observe se usa blockquote e aria-label.", 0.5, 6.7, 12.3, 0.4, size=13, color=GREEN)

# ── SLIDE 7: PILAR 3 AGENTS ──────────────────────────────────
s = slide()
pilar_header(s, "3", "AGENTES", "Agents",
    "Especialista terceirizado",
    "Quando sua empresa precisa de auditoria fiscal, contrata um contador especialista.\nEle tem ferramentas proprias e entrega no formato que voce precisa.")
txb(s, "O que e:", 0.5, 2.65, 3, 0.35, size=14, bold=True, color=GREEN)
txb(s, "Agente com identidade, escopo e ferramentas proprias. Invocado com @nome.", 0.5, 3.05, 12.3, 0.4, size=15)
ainfo = [("Nome","@qa-academia"),("Missao","Revisor de qualidade do site"),("Especialidade","Acessibilidade, SEO, HTML semantico"),("Formato","VERMELHO Bloqueante | AMARELO Atencao | VERDE Sugestao")]
for i,(k,v) in enumerate(ainfo):
    box(s, 0.5+(i%2)*6.2, 3.6+(i//2)*0.75, 5.9, 0.65, CARD_BG)
    txb(s, k+":", 0.7+(i%2)*6.2, 3.63+(i//2)*0.75, 1.5, 0.3, size=11, bold=True, color=GREEN)
    txb(s, v, 2.1+(i%2)*6.2, 3.63+(i//2)*0.75, 4.3, 0.3, size=12)
txb(s, "Instructions = geram codigo certo  |  Agent @qa-academia = audita o que foi gerado", 0.5, 5.5, 12.3, 0.4, size=14, color=GRAY, italic=True)
txb(s, "Valide: invoque @qa-academia e peca para revisar src/sections/sobre.html.", 0.5, 6.7, 12.3, 0.4, size=13, color=GREEN)

# ── SLIDE 8: PILAR 4 MCP ─────────────────────────────────────
s = slide()
pilar_header(s, "4", "MCP", "MCP -- Model Context Protocol",
    "O garcom do restaurante",
    "Voce nao vai a cozinha buscar a comida. O garcom vai, traz o que voce pediu e entrega na mesa.\nO garcom nao cozinha -- ele conecta voce a cozinha. O MCP funciona exatamente assim.")
txb(s, "Servidores configurados:", 0.5, 2.65, 5, 0.35, size=14, bold=True, color=GREEN)
servers = [
    ("fetch",           "Busca qualquer URL publica em tempo real (docs, APIs, conteudo externo)"),
    ("filesystem",      "Acessa arquivos do projeto com busca semantica"),
    ("precos-fitcode",  "Servidor local com precos reais da Academia FitCode"),
]
for i,(srv,desc) in enumerate(servers):
    box(s, 0.5, 3.1+i*0.72, 12.3, 0.62, CARD_BG)
    txb(s, srv, 0.7, 3.13+i*0.72, 2.5, 0.3, size=13, bold=True, color=GREEN)
    txb(s, desc, 3.4, 3.13+i*0.72, 9.2, 0.3, size=13)
txb(s, "Fluxo: 'Qual e o preco de Natacao?' -> Copilot chama consultar_precos -> MCP busca -> Resposta real", 0.5, 5.4, 12.3, 0.4, size=13, color=GRAY, italic=True)
txb(s, "Valide: pergunte o preco do plano mensal de Natacao -- o Copilot usara a tool, nao vai inventar.", 0.5, 6.7, 12.3, 0.4, size=13, color=GREEN)

# ── SLIDE 9: PILAR 5 HOOKS ───────────────────────────────────
s = slide()
pilar_header(s, "5", "HOOKS", "Hooks",
    "O porteiro do predio",
    "Toda vez que alguem tenta sair com um pacote, o porteiro confere: tem autorizacao?\nSe nao, o pacote volta. Os Hooks sao esse porteiro para o codigo do Copilot.")
txb(s, "Eventos disponiveis:", 0.5, 2.65, 5, 0.35, size=14, bold=True, color=GREEN)
eventos = [
    ("sessionStart",        "Inicio da sessao",        "Log de inicio, aviso de politicas"),
    ("preToolUse [UNICO]",  "ANTES de executar tool",  "UNICO que pode BLOQUEAR a acao"),
    ("postToolUse",         "Apos executar tool",       "Validar resultado, rodar lint"),
    ("userPromptSubmitted", "Ao enviar prompt",         "Filtrar dados sensiveis"),
    ("errorOccurred",       "Quando ocorre erro",       "Alerta automatico + rollback"),
]
for i,(ev,quando,oque) in enumerate(eventos):
    bg3 = RGBColor(0x10,0x25,0x1a) if "UNICO" in ev else CARD_BG
    box(s, 0.5, 3.1+i*0.6, 12.3, 0.52, bg3)
    col3 = AMBER if "UNICO" in ev else GREEN
    txb(s, ev, 0.7, 3.13+i*0.6, 3.0, 0.28, size=11, bold=True, color=col3)
    txb(s, quando, 3.85, 3.13+i*0.6, 3.5, 0.28, size=11)
    txb(s, oque, 7.5, 3.13+i*0.6, 4.8, 0.28, size=11, color=GRAY)
txb(s, "Valide: execute o lint-guard.ps1 no arquivo src/test-hook.html com os erros intencionais.", 0.5, 6.7, 12.3, 0.4, size=13, color=GREEN)

# ── SLIDE 10: PILAR 6 KNOWLEDGE BASE ─────────────────────────
s = slide()
pilar_header(s, "6", "KNOWLEDGE", "Knowledge Base",
    "A biblioteca da empresa",
    "Toda empresa serio tem intranet: manuais de marca, tabelas de precos, procedimentos.\nUm funcionario nao sai inventando -- ele consulta a intranet.")
txb(s, "Arquivos da wiki:", 0.5, 2.65, 5, 0.35, size=14, bold=True, color=GREEN)
wiki = [
    ("identidade-visual.md",  "Paleta de cores, tipografia, logo, tom de voz"),
    ("conteudo-academia.md",  "Modalidades, beneficios, publico-alvo, frases proibidas"),
    ("suporte-sla.md",        "Politicas de suporte e atendimento ao aluno"),
]
for i,(f,d) in enumerate(wiki):
    box(s, 0.5, 3.1+i*0.7, 12.3, 0.6, CARD_BG)
    txb(s, "docs/wiki/"+f, 0.7, 3.13+i*0.7, 4.8, 0.3, size=12, bold=True, color=GREEN)
    txb(s, d, 5.7, 3.13+i*0.7, 6.8, 0.3, size=12)
txb(s, "Skills = procedimento especifico  |  Knowledge Base = indice de toda documentacao", 0.5, 5.3, 12.3, 0.35, size=14, color=GRAY)
txb(s, "Conteudo mudou? Atualiza o Markdown. O agente usa a informacao nova automaticamente.", 0.5, 5.75, 12.3, 0.4, size=13, color=GRAY, italic=True)
txb(s, "Valide: pergunte 'Qual e a cor primaria da marca?' sem indicar o arquivo.", 0.5, 6.7, 12.3, 0.4, size=13, color=GREEN)

# ── SLIDE 11: PILAR 7 CODING AGENT ───────────────────────────
s = slide()
pilar_header(s, "7", "CODING", "Coding Agent",
    "O funcionario autonomo",
    "Voce manda a tarefa por escrito, ele planeja, executa e te manda o resultado para aprovar.\nVoce nao ficou do lado dele -- ele trabalhou de forma autonoma.")
txb(s, "Fluxo completo:", 0.5, 2.65, 5, 0.35, size=14, bold=True, color=GREEN)
steps = [
    ("1 - Atribuir issue",  "Crie a issue no GitHub e atribua ao Copilot"),
    ("2 - Copilot planeja", "Le AGENTS.md, instructions e explora o repo"),
    ("3 - Codifica",        "Escreve HTML/CSS, roda lint e corrige falhas"),
    ("4 - Abre o PR",       "Draft PR com codigo e link do deploy de preview"),
]
for i,(etapa,desc) in enumerate(steps):
    box(s, 0.5+(i%2)*6.2, 3.1+(i//2)*0.85, 5.9, 0.72, CARD_BG)
    txb(s, etapa, 0.7+(i%2)*6.2, 3.13+(i//2)*0.85, 2.8, 0.32, size=12, bold=True, color=GREEN)
    txb(s, desc, 0.7+(i%2)*6.2, 3.45+(i//2)*0.85, 5.5, 0.3, size=11)
txb(s, "VS Code Agent Mode: voce presente, guiando | Coding Agent: Copilot trabalha sozinho, Issue -> PR", 0.5, 5.2, 12.3, 0.5, size=13, color=GRAY)
txb(s, "Valide: crie uma issue, atribua ao Copilot, aguarde o Draft PR automatico.", 0.5, 6.7, 12.3, 0.4, size=13, color=GREEN)

# ── SLIDE 12: PILAR 8 CLI ────────────────────────────────────
s = slide()
pilar_header(s, "8", "CLI", "CLI",
    "O Copilot fora da janela grafica",
    "Nem todo mundo usa interface grafica. Sysadmins, devops vivem no terminal.\nO CLI garante que o Copilot esta la tambem -- sem botoes, sem mouse, sem IDE.")
txb(s, "Dois modos:", 0.5, 2.65, 4, 0.35, size=14, bold=True, color=GREEN)
box(s, 0.5, 3.1, 5.9, 2.7, CARD_BG)
txb(s, "gh copilot", 0.7, 3.15, 5.5, 0.4, size=16, bold=True, color=GREEN)
txb(s, "Resposta rapida, sem contexto do projeto", 0.7, 3.6, 5.5, 0.3, size=12, color=GRAY)
txb(s, 'gh copilot suggest "como fazer deploy"\ngh copilot explain "git rebase -i HEAD~3"', 0.7, 3.95, 5.5, 0.7, size=12)
txb(s, "Otimo para comandos git e shell rapidos", 0.7, 4.72, 5.5, 0.3, size=11, color=GRAY, italic=True)
box(s, 6.9, 3.1, 5.9, 2.7, CARD_BG)
txb(s, "copilot CLI", 7.1, 3.15, 5.5, 0.4, size=16, bold=True, color=GREEN)
txb(s, "Agente completo com contexto do projeto", 7.1, 3.6, 5.5, 0.3, size=12, color=GRAY)
txb(s, "copilot  # inicia sessao\n> Verifique imagens sem alt\n> &Gere secao em background\n> /resume", 7.1, 3.95, 5.5, 0.7, size=12)
txb(s, "Otimo para tarefas de codigo no terminal", 7.1, 4.72, 5.5, 0.3, size=11, color=GRAY, italic=True)
txb(s, 'Valide: gh copilot suggest "como encontrar HTMLs sem meta description"', 0.5, 6.7, 12.3, 0.4, size=13, color=GREEN)

# ── SLIDE 13: EXECUTE OS BLOCOS ──────────────────────────────
s = slide()
txb(s, "Execute os Blocos", 0.5, 0.3, 10, 0.7, size=36, bold=True, color=GREEN)
divider(s, 1.1)
blocos = [
    ("Bloco 1","bloco1-sem-padrao",         "Agente sem governanca -- HTML sem a11y"),
    ("Bloco 2","bloco2-com-instructions",   "Instructions + Skills -- HTML acessivel"),
    ("Bloco 3","bloco3-mcp-fetch",          "MCP -- agente busca diretrizes WCAG ao vivo"),
    ("Bloco 4","bloco4-adicionar-animacao", "Hook -- lint dispara e agente corrige"),
    ("Bloco 5","@qa-academia no chat",      "Custom Agent revisor de acessibilidade"),
    ("Bloco 6","bloco6-knowledge-base",     "Knowledge Base -- wiki -> codigo"),
    ("Bloco 7","Atribuir issue ao Copilot", "Coding Agent -- Issue -> PR + deploy"),
    ("Bloco 8","gh copilot / copilot CLI",  "CLI"),
]
for i,(bloco,como,demonstra) in enumerate(blocos):
    bg4 = RGBColor(0x12,0x1a,0x10) if i%2==0 else CARD_BG
    box(s, 0.5, 1.3+i*0.66, 12.3, 0.58, bg4)
    txb(s, bloco,    0.7,  1.33+i*0.66, 1.4, 0.28, size=11, bold=True, color=GREEN)
    txb(s, como,     2.2,  1.33+i*0.66, 3.7, 0.28, size=11)
    txb(s, demonstra,6.1,  1.33+i*0.66, 6.3, 0.28, size=11, color=GRAY)
txb(s, "Os prompts aparecem no Copilot Chat  ( / para listar )", 0.5, 6.9, 12.3, 0.35, size=13, color=GRAY, italic=True)

# ── SLIDE 14: GERE SUA EVIDENCIA ─────────────────────────────
s = slide()
txb(s, "Gere sua Evidencia", 0.5, 0.3, 12, 0.7, size=36, bold=True)
divider(s, 1.1)
box(s, 0.5, 1.3, 12.3, 1.1, DARK)
txb(s, 'python setup/gerar_evidencia.py --nome "Seu Nome Completo" --turma front', 0.7, 1.5, 12.0, 0.5, size=18, bold=True, color=GREEN)
txb(s, "O HTML abrira automaticamente. Voce deve ver 8/8 pilares configurados.", 0.5, 2.7, 12.3, 0.5, size=18)
box(s, 0.5, 3.45, 12.3, 1.7, CARD_BG)
txb(s, "Checklist final:", 0.7, 3.58, 11.8, 0.35, size=14, bold=True, color=GREEN)
checks = ["8/8 pilares configurados","Site renderiza no navegador","Screenshot do card de evidencia tirado","PR do Coding Agent revisado"]
for i,c in enumerate(checks):
    txb(s, "  OK  "+c, 0.7, 3.98+i*0.27, 11.8, 0.25, size=13)
txb(s, "Tire o screenshot do card e envie como evidencia de conclusao.", 0.5, 5.4, 12.3, 0.4, size=16, color=GRAY, italic=True)
r2 = s.shapes.add_shape(1,Inches(0),Inches(6.1),Inches(13.33),Inches(1.4))
r2.fill.solid(); r2.fill.fore_color.rgb = GREEN; r2.line.color.rgb = GREEN
txb(s, "Bom workshop!", 0.5, 6.25, 12.3, 0.9, size=38, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

out = r"C:\Users\edilson.b.monteiro\projects\Workshop\Site_academia\docs\workshop-participante.pptx"
prs.save(out)
print("OK: " + out)
